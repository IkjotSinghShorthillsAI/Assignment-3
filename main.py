from abc import ABC, abstractmethod
import os
import pdfplumber
import docx
import pptx
import mysql.connector
import csv
from PIL import Image
from io import BytesIO
import json
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE
from docx.oxml.ns import qn
from docx.document import Document as DocxDocument
import fitz  # PyMuPDF for PDF and font extraction

# Abstract Class: FileLoader
class FileLoader(ABC):
    def __init__(self, file_path):
        self.file_path = file_path
        self.validate_file()

    @abstractmethod
    def validate_file(self):
        pass

    @abstractmethod
    def load_file(self):
        pass

# Concrete Class: PDFLoader
class PDFLoader(FileLoader):
    def validate_file(self):
        if not self.file_path.endswith('.pdf'):
            raise ValueError("Invalid file type. Expected a PDF file.")
    
    def load_file(self):
        return pdfplumber.open(self.file_path)

# Concrete Class: DOCXLoader
class DOCXLoader(FileLoader):
    def validate_file(self):
        if not self.file_path.endswith('.docx'):
            raise ValueError("Invalid file type. Expected a DOCX file.")
    
    def load_file(self):
        return docx.Document(self.file_path)

# Concrete Class: PPTLoader
class PPTLoader(FileLoader):
    def validate_file(self):
        if not self.file_path.endswith('.pptx'):
            raise ValueError("Invalid file type. Expected a PPTX file.")
    
    def load_file(self):
        return pptx.Presentation(self.file_path)



class DataExtractor:
    def __init__(self, loader):
        self.loader = loader.load_file()
        self.file_path = loader.file_path
        self.file_name = os.path.basename(self.file_path)
        self.metadata = self.get_metadata()

    def get_metadata(self):
        file_stats = os.stat(self.file_path)
        return {
            "file_size": file_stats.st_size,
            "creation_time": file_stats.st_ctime,
            "modification_time": file_stats.st_mtime
        }

    def extract_text(self):
        results = []

        def is_heading(font_size, bold):
            try:
                fs = float(font_size)
            except (ValueError, TypeError):
                fs = 0
            return bold or fs > 12  # adjust threshold as needed

        # PDF Extraction using fitz
        if self.file_path.endswith('.pdf'):
            doc = fitz.open(self.file_path)
            for i, page in enumerate(doc):
                blocks = page.get_text("dict")["blocks"]
                for block in blocks:
                    if block["type"] == 0:  # text block
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span.get('text', '').strip()
                                if not text:
                                    continue
                                font_name = span.get("font", "Default")
                                font_size = span.get("size", 0)
                                bold = "Bold" in font_name
                                italic = False  # fitz does not provide italic info directly
                                data_type = "heading" if is_heading(font_size, bold) else "text"
                                results.append((i + 1, text, data_type, font_name, font_size, bold, italic))
            return results

        # DOCX Extraction using python-docx
        elif self.file_path.endswith('.docx'):
            doc = self.loader  # already loaded as a Document
            for para in doc.paragraphs:
                for run in para.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    font_name = run.font.name if run.font.name else "Default"
                    font_size = run.font.size.pt if run.font.size else 0
                    bold = run.bold if run.bold is not None else False
                    italic = run.italic if run.italic is not None else False
                    data_type = "heading" if is_heading(font_size, bold) else "text"
                    results.append((1, text, data_type, font_name, font_size, bold, italic))
            return results

        # PPTX Extraction using python-pptx
        elif self.file_path.endswith('.pptx'):
            pres = self.loader  # already loaded as a Presentation
            for i, slide in enumerate(pres.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                text = run.text.strip()
                                if not text:
                                    continue
                                font = run.font
                                font_name = font.name if font.name else "Default"
                                font_size = font.size.pt if font.size else 0
                                bold = font.bold if font.bold is not None else False
                                italic = font.italic if font.italic is not None else False
                                data_type = "heading" if is_heading(font_size, bold) else "text"
                                results.append((i + 1, text, data_type, font_name, font_size, bold, italic))
            return results

        return results
    def extract_links(self):
            links = []
            
            if isinstance(self.loader, pdfplumber.PDF):
                for i, page in enumerate(self.loader.pages):
                    if hasattr(page, 'annots') and page.annots:
                        for annot in page.annots:
                            uri = annot.get("uri")
                            if uri:
                                links.append((i + 1, uri))
            
            elif isinstance(self.loader, DocxDocument):
                for para in self.loader.paragraphs:
                    for rel in self.loader.part.rels.values():
                        if "hyperlink" in rel.reltype:
                            if para.text and rel.target_ref in para.text:
                                links.append((1, rel.target_ref))
            
            elif hasattr(self.loader, "slides"):  # Check if it's a Presentation object
                for slide_num, slide in enumerate(self.loader.slides):
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.hyperlink and run.hyperlink.address:
                                        links.append((slide_num + 1, run.hyperlink.address))
            
            return links

    def extract_images(self):
        images = []
        
        if self.file_path.endswith(".pdf"):
            doc = fitz.open(self.file_path)
            for page_num in range(len(doc)):
                for img_index, img in enumerate(doc[page_num].get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    img_obj = Image.open(BytesIO(image_bytes))
                    images.append((page_num + 1, "PNG", img_obj.size, img_obj))
        
        elif self.file_path.endswith(".docx"):
            for i, rel in enumerate(self.loader.part.rels):
                if "image" in self.loader.part.rels[rel].target_ref:
                    image_data = self.loader.part.rels[rel].target_part.blob
                    img_obj = Image.open(BytesIO(image_data))
                    images.append((i + 1, "PNG", img_obj.size, img_obj))
        
        elif self.file_path.endswith(".pptx"):
            for slide_num, slide in enumerate(self.loader.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture shape type
                        image = shape.image
                        img_obj = Image.open(BytesIO(image.blob))
                        images.append((slide_num + 1, "PNG", img_obj.size, img_obj))
        
        return images

    def extract_tables(self):
        tables = []
        if isinstance(self.loader, pdfplumber.PDF):
            for i, page in enumerate(self.loader.pages):
                extracted_tables = page.extract_tables()
                for table in extracted_tables:
                    tables.append((i + 1, len(table), len(table[0]) if table else 0, table))
        
        elif isinstance(self.loader, DocxDocument):
            for i, table in enumerate(self.loader.tables):
                extracted_table = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                tables.append((i + 1, len(extracted_table), len(extracted_table[0]) if extracted_table else 0, extracted_table))
        
        elif hasattr(self.loader, "slides"):  # Check if it's a Presentation object
            for slide_num, slide in enumerate(self.loader.slides):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        extracted_table = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                        tables.append((slide_num + 1, len(extracted_table), len(extracted_table[0]) if extracted_table else 0, extracted_table))
        
        return tables
    
# Abstract Class: Storage
class Storage(ABC):
    def __init__(self, extractor):
        self.extractor = extractor

    @abstractmethod
    def save_data(self):
        pass

# Concrete Class: FileStorage
class FileStorage(Storage):
    def save_data(self, output_dir):
        os.makedirs(output_dir, exist_ok=True)

        # Save text data with metadata
        with open(os.path.join(output_dir, "text_data.txt"), "w", encoding="utf-8") as f:
            for page_num, text, data_type, font_name, font_size, bold, italic in self.extractor.extract_text():
                f.write(
                    f"Page {page_num} - {data_type}: {text} "
                    f"(Font: {font_name}, Size: {font_size}, Bold: {bold}, Italic: {italic})\n"
                )

        # Save extracted tables
        with open(os.path.join(output_dir, "tables.csv"), "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for page_num, rows, cols, table in self.extractor.extract_tables():
                writer.writerow([f"Page {page_num} ({rows}x{cols})"])
                writer.writerows(table)

        # Save images
        for i, img_format, size, img in self.extractor.extract_images():
            image_path = os.path.join(output_dir, f"image_{i}.png")
            img.save(image_path, img_format)


# Concrete Class: SQLStorage with font metadata support
class SQLStorage(Storage):
    def __init__(self, extractor, host, user, password, database):
        super().__init__(extractor)
        self.host = host
        self.user = user
        self.password = password
        self.database = database
        self._ensure_database_exists()

    def _ensure_database_exists(self):
        conn = mysql.connector.connect(host=self.host, user=self.user, password=self.password)
        cursor = conn.cursor()
        cursor.execute(f"CREATE DATABASE IF NOT EXISTS {self.database}")
        conn.close()

    def save_data(self):
        conn = mysql.connector.connect(host=self.host, user=self.user, password=self.password, database=self.database)
        cursor = conn.cursor()

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS extracted_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                file_name VARCHAR(255),
                file_size BIGINT,
                creation_time DATETIME,
                modification_time DATETIME,
                page_number INT,
                data_type VARCHAR(50),
                content TEXT,
                font_name VARCHAR(100),
                font_size VARCHAR(50),
                bold BOOLEAN,
                italic BOOLEAN
            )
        """)

        metadata = self.extractor.metadata
        file_name = self.extractor.file_name

        for page_num, text, data_type, font_name, font_size, bold, italic in self.extractor.extract_text():
            cursor.execute("""
                INSERT INTO extracted_data (
                    file_name, file_size, creation_time, modification_time,
                    page_number, data_type, content, font_name, font_size, bold, italic
                )
                VALUES (
                    %s, %s, FROM_UNIXTIME(%s), FROM_UNIXTIME(%s),
                    %s, %s, %s, %s, %s, %s, %s
                )
            """, (
                file_name, metadata["file_size"], metadata["creation_time"], metadata["modification_time"],
                page_num, data_type, text, font_name, font_size, bold, italic
            ))

        for page_num, url in self.extractor.extract_links():
            cursor.execute("""
                INSERT INTO extracted_data (
                    file_name, file_size, creation_time, modification_time,
                    page_number, data_type, content, font_name, font_size, bold, italic
                )
                VALUES (
                    %s, %s, FROM_UNIXTIME(%s), FROM_UNIXTIME(%s),
                    %s, %s, %s, %s, %s, %s, %s
                )
            """, (file_name, metadata["file_size"], metadata["creation_time"],
                  metadata["modification_time"], page_num, "link", url, "", 0, False, False))

        for page_num, rows, cols, table in self.extractor.extract_tables():
            table_json = json.dumps(table)
            cursor.execute("""
                INSERT INTO extracted_data (
                    file_name, file_size, creation_time, modification_time,
                    page_number, data_type, content, font_name, font_size, bold, italic
                )
                VALUES (
                    %s, %s, FROM_UNIXTIME(%s), FROM_UNIXTIME(%s),
                    %s, %s, %s, %s, %s, %s, %s
                )
            """, (file_name, metadata["file_size"], metadata["creation_time"],
                  metadata["modification_time"], page_num, "table", table_json, "", 0, False, False))

        conn.commit()
        conn.close()

# Example Usage
def main():
    pdf_loader = PDFLoader("/home/shtlp_0096/Desktop/coding/assignment_3_dev/media/test1.pdf")
    doc_loader = DOCXLoader("/home/shtlp_0096/Desktop/coding/assignment_3_dev/media/demo.docx")
    ppt_loader = PPTLoader("/home/shtlp_0096/Desktop/coding/assignment_3_dev/media/ppt_test.pptx")
    extractor_ppt = DataExtractor(ppt_loader)
    extractor_pdf = DataExtractor(pdf_loader)
    extractor_doc = DataExtractor(doc_loader)

    for extractor, file_type in zip([extractor_pdf, extractor_doc, extractor_ppt], ["pdf", "doc", "ppt"]):
        file_storage = FileStorage(extractor)
        file_storage.save_data(f"output_{file_type}")

        sql_storage = SQLStorage(extractor, host="localhost", user="root", password="shills123", database="document_data")
        sql_storage.save_data()

if __name__ == "__main__":
    main()